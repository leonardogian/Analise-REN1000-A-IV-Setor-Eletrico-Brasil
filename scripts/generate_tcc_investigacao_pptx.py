from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "output" / "apresentacao_tcc_investigacao_dados_analises.pptx"
IMAGES = ROOT / "docs" / "images"

# Palette: ocean + midnight
C_DARK = RGBColor(6, 36, 64)
C_MID = RGBColor(14, 90, 128)
C_ACCENT = RGBColor(32, 177, 169)
C_LIGHT = RGBColor(240, 248, 252)
C_TEXT = RGBColor(35, 44, 51)
C_MUTED = RGBColor(98, 115, 129)


PRE_POST_METRICS = [
    ("Taxa fora do prazo (pre 2022)", "3.600%"),
    ("Taxa fora do prazo (pos 2022)", "2.524%"),
    ("Variacao absoluta", "-1.076 p.p."),
    ("Compensacao pre 2022", "R$ 335,3 mi"),
    ("Compensacao pos 2022", "R$ 87,6 mi"),
    ("Delta compensacao", "R$ -247,7 mi"),
]

NEO_HIGHLIGHTS = [
    "Coelba: 44,13 -> 28,39 fora prazo por 100k UC-mes (2023->2025)",
    "Cosern: 36,13 -> 9,38 fora prazo por 100k UC-mes",
    "Elektro: 7,37 -> 2,00 fora prazo por 100k UC-mes",
    "Pernambuco: 42,17 -> 17,71 fora prazo por 100k UC-mes",
    "Brasilia: estabilizacao em torno de 21 fora prazo por 100k UC-mes",
]

LIMITS = [
    "Anos 2024 e 2025 ainda incompletos para inferencia regulatoria final",
    "Bases com granularidades distintas (serie longa anual vs detalhe mensal)",
    "Leitura causal exige controle de efeitos operacionais e economicos",
]

NEXT_STEPS = [
    "Atualizar mensalmente fato_transgressao_mensal_porte",
    "Publicar endpoint FastAPI dedicado para transgressoes por grupo economico",
    "Adicionar alertas de mudanca abrupta (taxa e compensacao por UC)",
    "Consolidar narrativa final do TCC com recorte REN 1000/2021",
]


def add_full_bg(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title(slide, text: str, dark: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.bold = True
    run.font.size = Pt(34)
    run.font.color.rgb = C_LIGHT if dark else C_DARK


def add_subtitle(slide, text: str, y: float, dark: bool = False, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(12.0), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.color.rgb = C_LIGHT if dark else C_MUTED


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, dark: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(20)
        p.font.color.rgb = C_LIGHT if dark else C_TEXT
        p.space_after = Pt(8)
        p.bullet = True


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_metrics_grid(slide, metrics: list[tuple[str, str]]) -> None:
    card_w, card_h = 3.85, 1.45
    start_x, start_y = 0.8, 1.6
    gap_x, gap_y = 0.35, 0.3

    for i, (label, value) in enumerate(metrics):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(229, 242, 247)
        card.line.color.rgb = RGBColor(190, 216, 227)

        val_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.18), Inches(card_w - 0.35), Inches(0.5))
        p1 = val_box.text_frame.paragraphs[0]
        p1.text = value
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.runs[0]
        r1.font.size = Pt(28)
        r1.font.bold = True
        r1.font.color.rgb = C_DARK

        lbl_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.76), Inches(card_w - 0.35), Inches(0.55))
        p2 = lbl_box.text_frame.paragraphs[0]
        p2.text = label
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.runs[0]
        r2.font.size = Pt(13)
        r2.font.color.rgb = C_TEXT


def add_footer(slide, text: str, dark: bool = False) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_MID if dark else RGBColor(225, 238, 245)
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.11), Inches(12.0), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.name = "Calibri"
    r.font.size = Pt(11)
    r.font.color.rgb = C_LIGHT if dark else C_MUTED


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "TCC ANEEL - Investigacao com dados e analises"
    prs.core_properties.author = "Projeto TCC REN 1000/2021"
    prs.core_properties.subject = "Transgressoes de prazo e compensacoes financeiras"

    blank = prs.slide_layouts[6]

    # 1. Capa
    s1 = prs.slides.add_slide(blank)
    add_full_bg(s1, C_DARK)
    add_title(s1, "Investigacao com Dados e Analises", dark=True)
    add_subtitle(s1, "TCC - Eficacia da REN 1000/2021 (ANEEL)", y=1.35, dark=True, size=22)
    add_bullets(
        s1,
        [
            "Foco: transgressoes de prazo de servicos comerciais",
            "Recorte: compensacoes financeiras pagas aos consumidores",
            "Base: series 2011-2023 + detalhe mensal 2023-2025",
        ],
        x=0.9,
        y=2.2,
        w=6.5,
        h=2.6,
        dark=True,
    )
    add_image_fit(s1, IMAGES / "dashboard_visao_geral.png", x=7.1, y=1.5, w=5.8, h=4.8)
    add_footer(s1, f"Gerado em {date.today().strftime('%d/%m/%Y')} | Fontes: ANEEL + pipeline interno", dark=True)

    # 2. Pergunta e metodo
    s2 = prs.slides.add_slide(blank)
    add_title(s2, "Pergunta de Pesquisa e Metodo")
    add_subtitle(s2, "Como mensurar mudancas apos a REN 1000/2021 com consistencia estatistica?", y=1.1)
    add_bullets(
        s2,
        [
            "Comparacao pre 2022 (REN 414) vs pos 2022 (REN 1000)",
            "Normalizacao por UC ativa para comparabilidade entre portes",
            "Painel analitico com fatos mensais e anuais (CSV/Parquet)",
            "Validacoes de schema e rastreabilidade no pipeline ETL",
        ],
        x=0.8,
        y=1.8,
        w=6.3,
        h=4.4,
    )
    add_image_fit(s2, IMAGES / "dashboard_diagnostico.png", x=7.2, y=1.6, w=5.5, h=4.9)
    add_footer(s2, "Arquitetura: ETL -> tabelas analiticas -> FastAPI -> dashboard web")

    # 3. Dados
    s3 = prs.slides.add_slide(blank)
    add_title(s3, "Base de Dados Investigada")
    add_bullets(
        s3,
        [
            "Qualidade Comercial: serie longa consolidada (2011-2023)",
            "INDGER Servicos: detalhe mensal/municipio (2023-2025)",
            "INDGER Dados Comerciais: UC ativa para normalizacao",
            "Metrica central: fora_prazo_por_100k_uc_mes",
            "Metrica financeira: compensacao_rs e compensacao_por_uc",
        ],
        x=0.8,
        y=1.5,
        w=6.4,
        h=4.8,
    )
    add_image_fit(s3, IMAGES / "dashboard_analise_regulatoria.png", x=7.0, y=1.35, w=5.9, h=5.2)
    add_footer(s3, "Dados tratados em data/processed/analysis e consumidos em JSON no frontend")

    # 4. Resultado agregado
    s4 = prs.slides.add_slide(blank)
    add_title(s4, "Resultados Agregados: Pre vs Pos REN 1000")
    add_subtitle(s4, "Indicadores consolidados em relatorio_aneel.md", y=1.05)
    add_metrics_grid(s4, PRE_POST_METRICS)
    add_footer(s4, "Leitura: queda agregada de taxa e de compensacao no recorte comparavel")

    # 5. Serie anual
    s5 = prs.slides.add_slide(blank)
    add_title(s5, "Serie Anual Consolidada (2011-2023)")
    add_bullets(
        s5,
        [
            "Pico de taxa fora do prazo em 2015 (4,906%)",
            "Reducao progressiva no fim da serie (2023 em 1,944%)",
            "Compensacao anual ainda relevante em valores absolutos",
            "Interpretacao depende do nivel de cobertura anual",
        ],
        x=0.8,
        y=1.45,
        w=5.7,
        h=4.5,
    )
    add_image_fit(s5, IMAGES / "dashboard_visao_geral.png", x=6.2, y=1.45, w=6.5, h=4.9)
    add_footer(s5, "Serie longa comparavel: foco em tendencia, nao em causalidade isolada")

    # 6. Neoenergia
    s6 = prs.slides.add_slide(blank)
    add_full_bg(s6, C_LIGHT)
    add_title(s6, "Benchmark Neoenergia (5 distribuidoras)")
    add_bullets(s6, NEO_HIGHLIGHTS, x=0.8, y=1.55, w=6.5, h=4.9)
    add_image_fit(s6, IMAGES / "dashboard_neoenergia.png", x=7.1, y=1.35, w=5.7, h=5.2)
    add_footer(s6, "Fonte: reports/neoenergia_diagnostico.md | periodo 2023-2025")

    # 7. Transgressoes e grupos
    s7 = prs.slides.add_slide(blank)
    add_title(s7, "Transgressoes por Grupos Economicos")
    add_subtitle(s7, "Camada holding -> franquias e filtro para compensacao rural", y=1.05)
    add_image_fit(s7, IMAGES / "dashboard_transgressoes_main.png", x=0.8, y=1.6, w=6.0, h=4.7)
    add_image_fit(s7, IMAGES / "dashboard_transgressoes_filtrado.png", x=6.95, y=1.6, w=5.55, h=4.7)
    add_footer(s7, "Dashboard dedicado para cruzar volume de transgressoes e R$ compensado")

    # 8. Riscos e limites
    s8 = prs.slides.add_slide(blank)
    add_title(s8, "Limites de Interpretacao")
    add_bullets(s8, LIMITS, x=0.9, y=1.75, w=12.0, h=2.7)
    add_bullets(
        s8,
        [
            "Decisao metodologica: tratar 2024-2025 como monitoramento, nao conclusao",
            "Comparacao principal para inferencia: pre 2022 vs 2022-2023",
        ],
        x=0.9,
        y=4.45,
        w=12.0,
        h=1.5,
    )
    add_footer(s8, "A validade do TCC depende de manter recorte temporal e controle de comparabilidade")

    # 9. Proximos passos
    s9 = prs.slides.add_slide(blank)
    add_full_bg(s9, C_DARK)
    add_title(s9, "Proximos Passos do Trabalho", dark=True)
    add_bullets(s9, NEXT_STEPS, x=0.9, y=1.8, w=7.2, h=3.8, dark=True)
    add_image_fit(s9, IMAGES / "dashboard_analise_regulatoria.png", x=7.6, y=1.7, w=4.9, h=3.9)
    add_subtitle(
        s9,
        "Entrega final: narrativa tecnica + evidencias visuais + contratos de dados validos.",
        y=6.15,
        dark=True,
        size=18,
    )
    add_footer(s9, "Projeto TCC REN 1000/2021 | Dashboard e API na porta 8050", dark=True)

    # 10. Encerramento
    s10 = prs.slides.add_slide(blank)
    add_full_bg(s10, C_MID)
    add_title(s10, "Resumo Executivo", dark=True)
    add_bullets(
        s10,
        [
            "Ha evidencia de reducao agregada na taxa fora do prazo apos 2022",
            "A compensacao financeira segue relevante e heterogenea por distribuidora",
            "Normalizacao por UC ativa foi essencial para comparacao justa",
            "Painel analitico permite acompanhamento continuo da regulacao",
        ],
        x=0.9,
        y=1.9,
        w=11.7,
        h=3.5,
        dark=True,
    )
    add_subtitle(s10, "Obrigado.", y=6.05, dark=True, size=28)
    add_footer(s10, "Contato do projeto: repositorio TCC_leo_main", dark=True)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(path)
